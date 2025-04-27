import streamlit as st
import os
from deep_translator import GoogleTranslator
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
from io import BytesIO

# 🖼️ Background style
page_bg_img = '''
<style>
.stApp {
background-image: url("https://images.unsplash.com/photo-1503676260728-1c00da094a0b");
background-size: cover;
background-position: center;
}

.glass-card {
background: rgba(255, 255, 255, 0.2);
backdrop-filter: blur(10px);
border-radius: 20px;
padding: 30px;
margin: 50px auto;
max-width: 700px;
box-shadow: 0 8px 32px rgba(0, 0, 0, 0.2);
}
</style>
'''

st.markdown(page_bg_img, unsafe_allow_html=True)

def translate_text(text, target_language="en"):
    translator = GoogleTranslator(target=target_language)
    return translator.translate(text=text)

def translate_docx(uploaded_file, target_language="en"):
    doc = Document(uploaded_file)
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            paragraph.text = translate_text(paragraph.text, target_language)
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def translate_pdf(uploaded_file, target_language="en"):
    pdf_in = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    new_pdf = fitz.open()
    for page in pdf_in:
        text = page.get_text()
        translated_text = translate_text(text, target_language)
        new_page = new_pdf.new_page(width=page.rect.width, height=page.rect.height)
        new_page.insert_text((50, 50), translated_text)
    output = BytesIO()
    new_pdf.save(output)
    output.seek(0)
    return output

def translate_pptx(uploaded_file, target_language="en"):
    prs = Presentation(uploaded_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            run.text = translate_text(run.text, target_language)
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

def main():
    with st.container():
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        st.title("📄🌍 Sənəd Tərcümə Agenti")

        uploaded_file = st.file_uploader("DOCX, PDF və ya PPTX faylı yükləyin", type=["docx", "pdf", "pptx"])

        languages = {
            "İngilis": "en",
            "Fransız": "fr",
            "Alman": "de",
            "İspan": "es",
            "Ərəb": "ar",
            "Sadələşdirilmiş Çin": "zh-CN",
            "Rus": "ru",
            "Yapon": "ja",
            "İtalyan": "it",
            "Türk": "tr",
            "Azərbaycan": "az"
        }

        selected_language = st.selectbox("Tərcümə dili seçin:", list(languages.keys()))
        target_language = languages[selected_language]

        if uploaded_file:
            file_ext = os.path.splitext(uploaded_file.name)[-1].lower()
            st.info(f"{file_ext.upper()[1:]} fayl aşkarlandı!")

            if st.button("Tərcümə et"):
                progress_bar = st.progress(0)
                progress_bar.progress(30)

                if file_ext == ".docx":
                    output = translate_docx(uploaded_file, target_language)
                    st.download_button("Tərcümə olunmuş DOCX-i yüklə", output, file_name=f"tercume_{uploaded_file.name}")

                elif file_ext == ".pdf":
                    output = translate_pdf(uploaded_file, target_language)
                    st.download_button("Tərcümə olunmuş PDF-i yüklə", output, file_name=f"tercume_{uploaded_file.name}")

                elif file_ext == ".pptx":
                    output = translate_pptx(uploaded_file, target_language)
                    st.download_button("Tərcümə olunmuş PPTX-i yüklə", output, file_name=f"tercume_{uploaded_file.name}")

                else:
                    st.error("Dəstəklənməyən fayl tipi!")

                progress_bar.progress(100)
                st.success("Tərcümə tamamlandı!")

        st.markdown('</div>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()
